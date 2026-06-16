# -*- coding: utf-8 -*-
"""
report-pptx : reports/{종목명}.md  ->  reports/pptx/{종목명}.pptx

투자증권 리서치센터 톤 + 맥킨지 컨설팅 덱 스타일의 고정 양식 PPTX 생성기.

맥킨지 스타일 요소:
  - 아이브로우(eyebrow) 영문 라벨 + 한글 제목 (좌측 정렬)
  - 제목 아래 풀폭 헤어라인 + 좌측 짧은 오렌지 액센트
  - 하단 출처 각주(Source) + 페이지 번호 + footer 라인
  - 세로 괘선 없는 미니멀 표(헤더 네이비, 본문 가로 헤어라인만), 숫자 우측 정렬
  - 넉넉한 여백

입력 마크다운은 아래 섹션 헤더(## ...)로 구획을 나눈다. 헤더 텍스트에
키워드가 포함되면 해당 슬라이드로 매핑된다(부분 일치, 한글 기준).

  ## 개요          -> 종목 개요
  ## 재무          -> 재무 요약 (표 + 최근 3개년)
  ## 가격 / 추세   -> 가격/추세
  ## 뉴스 / 심리   -> 뉴스·심리
  ## 리스크        -> 리스크
  ## 종합          -> 한 줄 종합

표는 표준 마크다운 파이프 표(| ... |)로 작성한다. 헤더행 다음의
구분선(---) 행은 자동 무시된다.

사용:
  python build_pptx.py "삼성전자"
  python build_pptx.py --in reports/삼성전자.md --out reports/pptx/삼성전자.pptx
  python build_pptx.py "삼성전자" --date 2026-06-16
"""
import os
import re
import sys
import argparse
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 디자인 토큰
NAVY = RGBColor(0x1F, 0x2A, 0x44)       # 제목 / 표 헤더 배경
ORANGE = RGBColor(0xF3, 0x73, 0x21)     # 한화 오렌지 — 강조 전용(절제)
BODY_GRAY = RGBColor(0x3C, 0x40, 0x48)  # 본문 텍스트
SLATE = RGBColor(0x8A, 0x94, 0xA6)      # 아이브로우 / 각주 / 페이지번호
HAIRLINE = RGBColor(0xC9, 0xCE, 0xD6)   # 얇은 구분선 / 표 본문 괘선
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"                       # 한글 폰트 고정 — 깨짐 방지

SLIDE_W = Inches(13.333)                # 16:9
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.75)

EYEBROW_Y = Inches(0.52)
TITLE_Y = Inches(0.80)
RULE_Y = Inches(1.52)
CONTENT_TOP = Inches(1.88)
CONTENT_BOTTOM = Inches(6.72)           # 이 아래로는 본문/표를 그리지 않음(잘림 방지)
FOOTER_RULE_Y = Inches(6.92)
FOOTER_TEXT_Y = Inches(6.98)

# (정규화키, 슬라이드 제목, 영문 아이브로우, 각주 자료문구)
SLIDE_ORDER = [
    ("개요", "종목 개요", "COMPANY OVERVIEW",
     "자료: DART 전자공시"),
    ("재무", "재무 요약 (최근 3개년)", "FINANCIALS",
     "자료: DART 전자공시 (연결 사업보고서)"),
    ("가격", "가격 / 추세", "PRICE & TREND",
     "자료: FinanceDataReader (일별·지연 데이터)"),
    ("뉴스", "뉴스 · 심리", "NEWS & SENTIMENT",
     "자료: 언론 보도 종합"),
    ("리스크", "리스크", "KEY RISKS",
     "자료: DART · FinanceDataReader · 언론 종합"),
    ("종합", "한 줄 종합", "KEY TAKEAWAY",
     "자료: 본 리포트 종합"),
]
SECTION_ALIASES = {
    "개요": ["개요", "기업개요", "종목개요"],
    "재무": ["재무", "실적", "재무요약"],
    "가격": ["가격", "추세", "차트", "주가"],
    "뉴스": ["뉴스", "심리", "센티", "이슈"],
    "리스크": ["리스크", "위험"],
    "종합": ["종합", "결론", "한줄", "의견"],
}

_NUMERIC_RE = re.compile(r"^[\(\)\-+]?\d[\d,\.]*\s*(%|원|조|억|배|pt|x)?\)?$")


# ---------------------------------------------------------------- 마크다운 파싱
def parse_markdown(text):
    sections = {}
    cur_key = None
    buf = []

    def flush():
        if cur_key is not None:
            sections.setdefault(cur_key, "")
            sections[cur_key] += "\n".join(buf).strip() + "\n"

    for line in text.splitlines():
        m = re.match(r"^#{1,3}\s+(.*)$", line)
        if m and not line.startswith("#### "):
            flush()
            buf = []
            cur_key = _match_section(m.group(1))
        else:
            if cur_key is not None:
                buf.append(line)
    flush()
    return sections


def _match_section(title):
    norm = re.sub(r"\s+", "", title)
    for key, aliases in SECTION_ALIASES.items():
        for a in aliases:
            if a in norm:
                return key
    return None


def extract_table(block):
    rows = []
    for line in block.splitlines():
        s = line.strip()
        if s.startswith("|") and s.endswith("|"):
            cells = [c.strip() for c in s.strip("|").split("|")]
            if all(re.fullmatch(r":?-{2,}:?", c) for c in cells if c):
                continue
            rows.append(cells)
    return rows if len(rows) >= 1 else None


def extract_bullets(block):
    out = []
    for line in block.splitlines():
        s = line.strip()
        if not s or s.startswith("|"):
            continue
        s = re.sub(r"^[-*]\s+", "", s)
        s = re.sub(r"^\d+\.\s+", "", s)
        s = s.replace("**", "")
        out.append(s)
    return out


# ---------------------------------------------------------------- 저수준 헬퍼
def _set_run(run, size, color, bold=False, spacing=None):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rpr = run._r.get_or_add_rPr()
    if spacing is not None:
        rpr.set("spc", str(int(spacing * 100)))  # 1/100 pt
    ea = rpr.makeelement(qn("a:ea"), {})
    ea.set("typeface", FONT)
    rpr.append(ea)


def _rule(slide, x, y, w, h, color, shadow=False):
    shp = slide.shapes.add_shape(1, x, y, w, h)  # rectangle
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = shadow
    return shp


def _cell_border(cell, edges=("bottom",), color="C9CED6", w_pt=0.75):
    """세로 괘선 없는 맥킨지 표를 위해 셀 테두리를 개별 지정."""
    tcPr = cell._tc.get_or_add_tcPr()
    tagmap = {"left": "a:lnL", "right": "a:lnR", "top": "a:lnT", "bottom": "a:lnB"}
    fill = None
    for fc in ("a:solidFill", "a:noFill", "a:gradFill", "a:blipFill",
               "a:pattFill", "a:grpFill"):
        fill = tcPr.find(qn(fc))
        if fill is not None:
            break
    for e in ("left", "right", "top", "bottom"):
        if e not in edges:
            continue
        tag = tagmap[e]
        old = tcPr.find(qn(tag))
        if old is not None:
            tcPr.remove(old)
        ln = tcPr.makeelement(qn(tag), {
            "w": str(int(Pt(w_pt))), "cap": "flat", "cmpd": "sng", "algn": "ctr",
        })
        sf = ln.makeelement(qn("a:solidFill"), {})
        clr = sf.makeelement(qn("a:srgbClr"), {"val": color})
        sf.append(clr)
        ln.append(sf)
        if fill is not None:
            fill.addprevious(ln)
        else:
            tcPr.append(ln)


def _strip_table_style(table):
    """기본 표 스타일(밴딩·테두리) 제거 → 수동 포맷만 적용."""
    tblPr = table._tbl.tblPr
    sid = tblPr.find(qn("a:tableStyleId"))
    if sid is not None:
        tblPr.remove(sid)
    table.first_row = False
    table.horz_banding = False
    table.last_row = False
    table.first_col = False


# ---------------------------------------------------------------- 슬라이드 골격
def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_header(slide, title, eyebrow):
    # 아이브로우(영문, 자간 넓게)
    eb = slide.shapes.add_textbox(MARGIN, EYEBROW_Y, SLIDE_W - 2 * MARGIN, Inches(0.3))
    p = eb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = eyebrow.upper()
    _set_run(r, 10.5, ORANGE, bold=True, spacing=2.2)
    # 제목
    tb = slide.shapes.add_textbox(MARGIN, TITLE_Y, SLIDE_W - 2 * MARGIN, Inches(0.7))
    tb.text_frame.word_wrap = True
    p2 = tb.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = title
    _set_run(r2, 25, NAVY, bold=True)
    # 풀폭 헤어라인 + 좌측 오렌지 액센트
    _rule(slide, MARGIN, RULE_Y, SLIDE_W - 2 * MARGIN, Pt(0.75), HAIRLINE)
    _rule(slide, MARGIN, RULE_Y - Pt(0.6), Inches(1.1), Pt(2.4), ORANGE)


def add_footer(slide, source, page_no):
    _rule(slide, MARGIN, FOOTER_RULE_Y, SLIDE_W - 2 * MARGIN, Pt(0.5), HAIRLINE)
    fb = slide.shapes.add_textbox(MARGIN, FOOTER_TEXT_Y, Inches(9.5), Inches(0.35))
    p = fb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = source + "   ·   학습용 분석이며 투자 권유가 아닙니다."
    _set_run(r, 8.5, SLATE)
    pn = slide.shapes.add_textbox(SLIDE_W - MARGIN - Inches(1.2), FOOTER_TEXT_Y,
                                  Inches(1.2), Inches(0.35))
    pp = pn.text_frame.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    rr = pp.add_run()
    rr.text = str(page_no)
    _set_run(rr, 9, SLATE)


def add_bullets(slide, bullets, top=CONTENT_TOP, size=14, emphasize=False):
    if not bullets:
        return
    left = MARGIN
    width = SLIDE_W - 2 * MARGIN
    if emphasize:
        # 좌측 오렌지 액센트 바 + 들여쓰기 (KEY TAKEAWAY 강조)
        _rule(slide, MARGIN, top + Inches(0.02), Pt(3.5),
              Inches(0.55) * len(bullets), ORANGE)
        left = MARGIN + Inches(0.25)
        width = width - Inches(0.25)
    box = slide.shapes.add_textbox(left, top, width, CONTENT_BOTTOM - top)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10 if not emphasize else 8)
        p.line_spacing = 1.15
        marker = p.add_run()
        marker.text = "▪  "
        _set_run(marker, size, ORANGE if emphasize else NAVY, bold=True)
        body = p.add_run()
        body.text = b
        _set_run(body, size, NAVY if emphasize else BODY_GRAY,
                 bold=emphasize)


def add_table(slide, rows, top=CONTENT_TOP):
    if not rows:
        return
    ncols = max(len(r) for r in rows)
    rows = [r + [""] * (ncols - len(r)) for r in rows]

    avail_h = CONTENT_BOTTOM - top
    row_h = Inches(0.46)
    max_rows = max(2, int(avail_h / row_h))
    truncated = False
    if len(rows) > max_rows:
        rows = rows[: max_rows - 1]
        truncated = True

    nrows = len(rows) + (1 if truncated else 0)
    width = SLIDE_W - 2 * MARGIN
    shape = slide.shapes.add_table(nrows, ncols, MARGIN, top, width, row_h * nrows)
    tbl = shape.table
    _strip_table_style(tbl)

    for ci in range(ncols):
        tbl.columns[ci].width = Emu(int(width / ncols))

    for ri, row in enumerate(rows):
        is_header = ri == 0
        tbl.rows[ri].height = row_h
        for ci in range(ncols):
            cell = tbl.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if is_header else WHITE
            txt = row[ci] if ci < len(row) else ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            numeric = bool(_NUMERIC_RE.match(txt.strip()))
            if ci == 0:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.RIGHT if (numeric or is_header) else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = txt
            _set_run(r, 11.5 if is_header else 11,
                     WHITE if is_header else BODY_GRAY, bold=is_header)
            # 괘선: 헤더는 하단 네이비 굵게, 본문은 하단 헤어라인만 (세로선 없음)
            if is_header:
                _cell_border(cell, edges=("bottom",), color="1F2A44", w_pt=1.5)
            else:
                _cell_border(cell, edges=("bottom",), color="C9CED6", w_pt=0.75)

    if truncated:
        ri = len(rows)
        cell = tbl.cell(ri, 0)
        if ncols > 1:
            cell.merge(tbl.cell(ri, ncols - 1))
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = "… 행이 많아 일부 생략되었습니다 (원본 .md 참조)"
        _set_run(r, 9.5, ORANGE, bold=True)

    return shape


# ---------------------------------------------------------------- 슬라이드 빌드
def build_cover(prs, stock, made_date):
    s = add_blank_slide(prs)
    # 좌측 풀하이트 네이비 액센트 밴드(얇게) + 그 위 오렌지 짧은 마크
    _rule(s, 0, 0, Inches(0.18), SLIDE_H, NAVY)
    _rule(s, 0, Inches(2.55), Inches(0.18), Inches(1.4), ORANGE)

    # 아이브로우
    eb = s.shapes.add_textbox(Inches(1.0), Inches(2.55), SLIDE_W - Inches(1.7), Inches(0.4))
    p = eb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "EQUITY RESEARCH  ·  가치투자 관점"
    _set_run(r, 12, SLATE, bold=True, spacing=2.0)

    # 메인 타이틀
    tb = s.shapes.add_textbox(Inches(1.0), Inches(2.95), SLIDE_W - Inches(1.7), Inches(1.4))
    tb.text_frame.word_wrap = True
    p2 = tb.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = stock
    _set_run(r2, 48, NAVY, bold=True)
    p3 = tb.text_frame.add_paragraph()
    r3 = p3.add_run()
    r3.text = "기업 분석 리포트"
    _set_run(r3, 22, BODY_GRAY)

    # 헤어라인
    _rule(s, Inches(1.0), Inches(5.35), SLIDE_W - Inches(2.0), Pt(0.75), HAIRLINE)

    # 메타데이터
    meta = s.shapes.add_textbox(Inches(1.0), Inches(5.5), SLIDE_W - Inches(2.0), Inches(0.5))
    pm = meta.text_frame.paragraphs[0]
    rm = pm.add_run()
    rm.text = f"분석 기준일  {made_date}      |      작성  리서치센터"
    _set_run(rm, 14, BODY_GRAY)

    # 하단 고지
    disc = s.shapes.add_textbox(Inches(1.0), Inches(6.75), SLIDE_W - Inches(2.0), Inches(0.5))
    pd = disc.text_frame.paragraphs[0]
    rd = pd.add_run()
    rd.text = ("※ 본 자료는 학습용이며 투자 권유가 아닙니다. 매수·매도 의견을 담지 않으며, "
               "최종 판단은 투자자 본인에게 있습니다.")
    _set_run(rd, 10, SLATE)


def build_content_slide(prs, key, title, eyebrow, source, block, page_no):
    s = add_blank_slide(prs)
    add_header(s, title, eyebrow)
    block = block or ""
    table = extract_table(block)
    if table:
        add_table(s, table)
    else:
        bullets = extract_bullets(block)
        if not bullets:
            bullets = ["확인 불가 — 원본 리포트에 해당 섹션 내용이 없습니다."]
        emphasize = key == "종합"
        add_bullets(s, bullets, size=16 if emphasize else 14, emphasize=emphasize)
    add_footer(s, source, page_no)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("stock", nargs="?", help="종목명 (reports/{종목명}.md)")
    ap.add_argument("--in", dest="infile", help="입력 .md 경로")
    ap.add_argument("--out", dest="outfile", help="출력 .pptx 경로")
    ap.add_argument("--date", dest="made_date", help="작성일 YYYY-MM-DD")
    ap.add_argument("--root", default=".", help="프로젝트 루트")
    args = ap.parse_args()

    root = os.path.abspath(args.root)
    if args.infile:
        infile = os.path.abspath(args.infile)
        stock = os.path.splitext(os.path.basename(infile))[0]
    elif args.stock:
        stock = args.stock
        infile = os.path.join(root, "reports", f"{stock}.md")
    else:
        ap.error("종목명 또는 --in 이 필요합니다.")

    if not os.path.exists(infile):
        sys.exit(f"[오류] 입력 파일이 없습니다: {infile}")

    outfile = args.outfile or os.path.join(root, "reports", "pptx", f"{stock}.pptx")
    outfile = os.path.abspath(outfile)
    os.makedirs(os.path.dirname(outfile), exist_ok=True)

    made_date = args.made_date or date.today().isoformat()

    with open(infile, encoding="utf-8") as f:
        text = f.read()
    sections = parse_markdown(text)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs, stock, made_date)
    for i, (key, title, eyebrow, source) in enumerate(SLIDE_ORDER, start=1):
        build_content_slide(prs, key, title, eyebrow, source,
                             sections.get(key, ""), page_no=i)

    prs.save(outfile)
    print(f"[완료] {outfile}")
    print(f"  슬라이드 {len(prs.slides._sldIdLst)}장 (표지 + {len(SLIDE_ORDER)}섹션)")
    missing = [t for k, t, _e, _s in SLIDE_ORDER if not sections.get(k)]
    if missing:
        print("  [주의] 내용이 비어 '확인 불가'로 채운 섹션: " + ", ".join(missing))


if __name__ == "__main__":
    main()
